#!/usr/bin/env python3
import sys
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from io import BytesIO

def extract_text_from_pptx(pptx_path, output_path="pptx_text_output.txt", ocr_langs="eng+ara"):
    if not os.path.isfile(pptx_path):
        raise FileNotFoundError(f"File not found: {pptx_path}")

    prs = Presentation(pptx_path)
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        lines.append("="*40)
        lines.append(f"Slide {i}")
        lines.append("="*40)

        # 1) Text from shapes (text frames, tables, etc.)
        for shape in slide.shapes:
            # Text frames (titles, text boxes)
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    lines.append("[TextBox] " + text)

            # Tables
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    table = shape.table
                    for r in range(len(table.rows)):
                        row_text = []
                        for c in range(len(table.columns)):
                            cell_text = table.cell(r, c).text.strip()
                            row_text.append(cell_text)
                        if any(cell.strip() for cell in row_text):
                            lines.append("[Table Row] " + " | ".join(row_text))
                except Exception:
                    pass

        # 2) OCR on pictures
        for shape in slide.shapes:
            # Detect picture shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(BytesIO(image_bytes)).convert("RGB")
                    # Run OCR. `ocr_langs` default includes eng+ara; change if needed.
                    try:
                        ocr_text = pytesseract.image_to_string(img, lang=ocr_langs)
                    except pytesseract.pytesseract.TesseractError:
                        # fallback to default language if requested langs not installed
                        ocr_text = pytesseract.image_to_string(img)
                    if ocr_text and ocr_text.strip():
                        lines.append("[OCR->Image] " + ocr_text.strip())
                except Exception as e:
                    lines.append(f"[OCR->Image Error] {e}")

        # 3) Speaker notes (if present)
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append("[Notes] " + notes_text)
        except Exception:
            # some slides may not have notes_slide
            pass

        lines.append("")  # blank line between slides

    # Write to output file
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"Done — output saved to: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 extract_pptx_ocr.py <presentation.pptx> [output.txt] [ocr_langs]")
        print("Example: python3 extract_pptx_ocr.py your_presentation.pptx output.txt eng+ara")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) >= 3 else "pptx_text_output.txt"
    ocr_langs = sys.argv[3] if len(sys.argv) >= 4 else "eng+ara"
    extract_text_from_pptx(pptx_path, output_file, ocr_langs)
